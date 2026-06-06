"""
Market Signals pipeline — research, analyse, generate PPTX + MD, post to Slack.

Environment variables:
  GROQ_API_KEY       required
  SERPER_API_KEY     required
  SLACK_BOT_TOKEN    required unless DRY_RUN=true
  SLACK_CHANNEL_ID   required unless DRY_RUN=true
  GITHUB_TOKEN       required unless DRY_RUN=true (for committing MD report)
  DRY_RUN            set to 'true' to skip Slack + git commit (save artifacts locally)
  DAYS_BACK          optional, default 45 — research window in days
"""
import json
import os
import subprocess
import sys
import time
from datetime import date, timedelta
from pathlib import Path

import requests
from groq import Groq
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

GROQ_MODEL = "llama-3.3-70b-versatile"

ELLI_GREEN    = RGBColor(0x00, 0xD4, 0x8A)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xBB, 0xBB, 0xBB)
SEPARATOR     = RGBColor(0x3A, 0x3A, 0x5A)
TEXT_ON_GREEN = RGBColor(0x05, 0x05, 0x18)

GEO_FOCUS = "Europe OR Germany OR Deutschland OR EU"

FLEET_PRODUCTS = (
    "Fleet Charging Console, Charging Card, Home Charging Reimbursement, "
    "Charge & Fuel Card, LOGPAY Card 4U"
)
FLEET_COMPETITORS = (
    "DKV Mobility, UTA Edenred, Shell Recharge, BP Pulse, EnBW, Aral, Octopus Energy"
)
FLEET_USERS = "Fleet Managers, Company Car Drivers, HR / Finance / Sustainability"

SITE_PRODUCTS = (
    "Site Management Console, Standortanalyse, Planning & Installation, "
    "Hardware Migration, Maintenance & Operations"
)
SITE_COMPETITORS = (
    "Driivz, AMPECO, Monta, Virta, ChargePoint, E.On, EnBW, "
    "The Mobility House, Moon"
)
SITE_USERS = "Commercial Property Owners, Dealers, Facilities Managers"

# Primary competitors get dedicated queries in addition to category queries.
FLEET_COMPETITOR_QUERIES = [
    ("DKV Mobility fleet EV charging card Europe {m}", "Competitor card and MSP moves"),
    ("UTA Edenred fleet EV charging card Europe {m}", "Competitor card and MSP moves"),
    ("Shell Recharge BP Pulse fleet charging Europe {m}", "Competitor card and MSP moves"),
    ("EnBW Aral fleet EV charging card Germany {m}", "Competitor card and MSP moves"),
    ("Octopus Energy fleet charging Europe {m}", "Competitor card and MSP moves"),
]

SITE_COMPETITOR_QUERIES = [
    ("Driivz AMPECO CPMS EV charging software Europe {m}", "CPMS platform competition and consolidation"),
    ("Monta Virta charging management platform Europe {m}", "CPMS platform competition and consolidation"),
    ("ChargePoint E.On EnBW charging site management Europe {m}", "CPMS platform competition and consolidation"),
    ("The Mobility House ChargePilot EV site management Europe {m}", "CPMS platform competition and consolidation"),
    ("Moon EV charging platform Europe {m}", "CPMS platform competition and consolidation"),
]

FLEET_CATEGORIES = [
    ("Pricing model and subscription evolution",
     f"fleet EV charging card fee kWh tariff pricing {GEO_FOCUS} {{m}}"),
    ("Network access and roaming coverage",
     f"eRoaming IONITY fleet charging network HPC partnership {GEO_FOCUS} {{m}}"),
    ("Mixed-fleet and ICE-to-EV transition",
     f"fuel card electrification ICE EV combined billing fleet {GEO_FOCUS} {{m}}"),
    ("Home charging reimbursement and compliance",
     f"home charging reimbursement kWh mandate wallbox fleet {GEO_FOCUS} {{m}}"),
    ("Plug&Charge and authentication technology",
     f"ISO 15118 Plug&Charge RFID fleet authentication Europe {{m}}"),
    ("Mobile payment and in-app fuelling",
     f"mobile app payment fuelling fleet Germany Europe {{m}}"),
    ("Employee benefit and non-cash incentive regulation",
     f"fleet benefit card Dienstwagen Steuervorteil EV charging Germany {{m}}"),
    ("Data, reporting and CO2 transparency",
     f"CSRD fleet Scope 3 CO2 EV reporting Fuhrpark Germany Europe {{m}}"),
    ("Fleet software consolidation and platform convergence",
     f"CPMS fleet card integration unified platform merger Europe {{m}}"),
    ("Commercial vehicle and HGV electrification",
     f"LKW Nutzfahrzeuge electric HGV LCV fleet charging Germany Europe {{m}}"),
    ("Competitor card and MSP moves",
     f"DKV UTA Shell BP EnBW Aral Fuhrpark EV Ladekarte news {GEO_FOCUS} {{m}}"),
    ("Regulatory and compliance obligations",
     f"AFIR fleet electrification regulation Germany EU {{m}}"),
    ("Customer pain points",
     f"fleet manager EV charging Fuhrpark problems Germany Europe {{m}}"),
]

SITE_CATEGORIES = [
    ("CPMS platform competition and consolidation",
     f"Driivz Ampeco Monta Virta ChargePoint CPMS Europe {{m}}"),
    ("Energy management and grid integration",
     f"EV charging Energiemanagement V2G demand response site Germany Europe {{m}}"),
    ("Hardware ecosystem and multi-vendor compatibility",
     f"OCPP 2.0.1 Alpitronic Compleo Mennekes charger integration Europe {{m}}"),
    ("Site planning and installation market",
     f"EV charging site Planung Netzanschluss permit cost Germany Europe {{m}}"),
    ("Public charging monetisation",
     f"CPO Ladesäule ad-hoc payment AFIR eRoaming revenue Europe {{m}}"),
    ("Building and property regulation",
     f"EPBD Gebäudeenergiegesetz pre-cabling building EV charging Germany Europe {{m}}"),
    ("CPO data and reporting obligations",
     f"DATEX II CPO Bundesnetzagentur data compliance Germany {{m}}"),
    ("Maintenance, uptime and operations",
     f"EV charging Verfügbarkeit remote diagnostics maintenance Germany Europe {{m}}"),
    ("AI and software-driven site operations",
     f"AI predictive maintenance EV charging autonomous network Europe {{m}}"),
    ("Segment-specific site demand",
     f"retail Gewerbe depot parking EV charging Germany Europe {{m}}"),
    ("Fleet and site convergence",
     f"fleet site unified EV charging console integration Europe {{m}}"),
    ("Sustainability and carbon reporting at site level",
     f"CO2 kWh renewable CSRD site charging Nachhaltigkeit Germany {{m}}"),
]

# ---------------------------------------------------------------------------
# Date helpers
# ---------------------------------------------------------------------------

MONTH_NAMES = [
    "", "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December",
]


def compute_date_range():
    days_back = int(os.getenv("DAYS_BACK", "45"))
    end = date.today()
    start = end - timedelta(days=days_back)

    current_month = MONTH_NAMES[end.month]
    prior_month = MONTH_NAMES[start.month]
    year = end.year

    if prior_month == current_month:
        month_range = f"{current_month} {year}"
    else:
        month_range = f"{prior_month} {year} OR {current_month} {year}"

    label = end.strftime("%Y-%m")
    generated_at = end.strftime("%B %Y")
    period_start = start.strftime("%-d %B %Y")
    period_end = end.strftime("%-d %B %Y")

    return {
        "start": start,
        "end": end,
        "month_range": month_range,
        "label": label,
        "generated_at": generated_at,
        "period_start": period_start,
        "period_end": period_end,
    }


# ---------------------------------------------------------------------------
# Phase 1 — Research
# ---------------------------------------------------------------------------

def search_serper(query, max_results=6):
    """Search via Serper.dev Google News. Free tier: 2500 queries/month."""
    try:
        resp = requests.post(
            "https://google.serper.dev/news",
            headers={
                "X-API-KEY": os.environ["SERPER_API_KEY"],
                "Content-Type": "application/json",
            },
            json={"q": query, "num": max_results, "gl": "de", "hl": "en"},
            timeout=15,
        )
        resp.raise_for_status()
        results = []
        for r in resp.json().get("news", []):
            results.append({
                "title": r.get("title", ""),
                "snippet": (r.get("snippet", "") or "")[:500],
                "url": r.get("link", ""),
                "source": r.get("source", ""),
                "date": r.get("date", ""),
            })
        return results
    except Exception as e:
        print(f"  search warning [{query[:50]}]: {e}", file=sys.stderr)
        return []


def run_searches(categories, competitor_queries, month_range):
    """
    Run all category searches plus dedicated competitor queries.
    Returns {category: [{idx,title,snippet,url,source,date}]}
    """
    results = {}

    # Category queries
    for i, (category, query_template) in enumerate(categories, 1):
        query = query_template.replace("{m}", month_range)
        hits = search_serper(query)
        indexed = [{"idx": j + 1, **h} for j, h in enumerate(hits)]
        # Merge into category bucket
        results.setdefault(category, [])
        existing_count = len(results[category])
        for h in indexed:
            h["idx"] = existing_count + h["idx"]
            results[category].append(h)
        print(f"  [{i}/{len(categories)}] {category[:50]}: {len(hits)} results")

    # Competitor-specific queries — merge into their category bucket
    for query_template, target_category in competitor_queries:
        query = query_template.replace("{m}", month_range)
        hits = search_serper(query, max_results=4)
        if hits:
            results.setdefault(target_category, [])
            existing_count = len(results[target_category])
            for j, h in enumerate(hits):
                results[target_category].append({"idx": existing_count + j + 1, **h})

    return results


# ---------------------------------------------------------------------------
# Phase 2 — Extract signals (Groq × batched calls)
# ---------------------------------------------------------------------------

_groq_client = None


def groq_client():
    global _groq_client
    if _groq_client is None:
        _groq_client = Groq(api_key=os.environ["GROQ_API_KEY"])
    return _groq_client


def call_groq(system_prompt, user_prompt, max_tokens=3000):
    """Call Groq in JSON mode; return parsed dict."""
    resp = groq_client().chat.completions.create(
        model=GROQ_MODEL,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.3,
        max_tokens=max_tokens,
    )
    return json.loads(resp.choices[0].message.content)


def _format_snippets_for_prompt(category, hits):
    lines = [f"--- {category} ---"]
    for h in hits:
        source_info = f" ({h.get('source', '')} {h.get('date', '')})" if h.get("source") else ""
        lines.append(f"[{h['idx']}] {h['title']}{source_info}")
        lines.append(f"URL: {h['url']}")
        lines.append(h["snippet"])
        lines.append("")
    return "\n".join(lines)


EXTRACT_SYSTEM = (
    "You extract concrete market signals from EV charging industry search results "
    "focused on the European market (Germany, Austria, Spain, Italy, France and EU). "
    "Be factual and specific — only report what the search results actually contain. "
    "Include company names, countries, and specific details where available. "
    "Do not invent signals not supported by the text."
)


def extract_signals(results_by_category, batch_size=4):
    """Extract concrete signals per category; returns {category: [{signal,sources}]}"""
    categories = list(results_by_category.keys())
    batches = [categories[i:i + batch_size] for i in range(0, len(categories), batch_size)]
    extracted = {}

    for b_idx, batch in enumerate(batches, 1):
        print(f"  extract batch {b_idx}/{len(batches)} ({len(batch)} categories)...")
        snippets_text = "\n".join(
            _format_snippets_for_prompt(cat, results_by_category[cat])
            for cat in batch
        )
        schema_example = {cat: [{"signal": "...", "sources": [1]}] for cat in batch}
        user_prompt = (
            f"Extract 2–4 concrete signals for each category: {', '.join(batch)}.\n\n"
            f"A signal = a specific product launch, pricing move, regulatory update, "
            f"partnership, competitor action, or customer pain point. "
            f"Include company names and countries.\n\n"
            f"Return JSON:\n{json.dumps(schema_example, indent=2)}\n\n"
            f"Search results:\n{snippets_text}"
        )
        try:
            data = call_groq(EXTRACT_SYSTEM, user_prompt)
            for cat in batch:
                extracted[cat] = data.get(cat, [])
        except Exception as e:
            print(f"  extract error batch {b_idx}: {e}", file=sys.stderr)
            for cat in batch:
                extracted[cat] = []
        time.sleep(1.0)

    return extracted


def build_url_map(results_by_category):
    url_map = {}
    for cat, hits in results_by_category.items():
        for h in hits:
            url_map[(cat, h["idx"])] = h["url"]
    return url_map


# ---------------------------------------------------------------------------
# Phase 3 — Synthesise report (Groq × 2 calls)
# ---------------------------------------------------------------------------

SYNTH_SYSTEM_TEMPLATE = """\
You are a senior market research analyst for Elli, a Volkswagen Group EV charging company \
operating primarily in Germany and Europe.

{area}
Products: {products}
Key B2B Competitors: {competitors}
Primary Users: {users}
Geographic focus: Germany (primary), Austria, Spain, Italy, France, EU

Analyse the provided market signals (past 30–45 days) and produce a detailed \
monthly Market Signals report for {area}. \
Write with the depth of a professional market research report — each item should \
explain WHAT the signal is, WHAT IT MEANS for the market, and what the implication \
is for Elli. Be specific: name companies, geographies, and figures where available.

Return ONLY valid JSON with this EXACT schema (no extra keys):
{{
  "trends": [
    {{
      "title": "Short headline (6–10 words)",
      "signal": "What is specifically happening — 3–4 sentences with company names, countries, and concrete details",
      "meaning": "Why this matters for the European EV charging market — 2–3 sentences on market direction",
      "implication": "What Elli should monitor, consider or act on — 1–2 sentences",
      "sources": [1, 2]
    }}
  ],
  "competitor_moves": [
    {{
      "company": "Competitor name",
      "action": "One-line summary of what they did",
      "detail": "2–3 sentences: context, geography, and why it matters competitively",
      "sources": [1]
    }}
  ],
  "unmet_needs": [
    {{
      "need": "The unmet need in one line",
      "evidence": "What the research shows — 2–3 sentences",
      "sources": [1]
    }}
  ],
  "active_regulations": [
    {{
      "regulation": "Regulation or policy name",
      "status": "Current implementation status",
      "impact": "Impact on fleet/site charging in Europe — 2–3 sentences",
      "sources": [1]
    }}
  ],
  "strategic_implications": [
    {{
      "implication": "Strategic recommendation or watch item for Elli",
      "rationale": "Why now and what the risk/opportunity is — 2–3 sentences"
    }}
  ],
  "open_questions": [
    {{"text": "A specific question for PM follow-up or further research"}}
  ]
}}

Produce exactly 5 trends. 3–5 items per other section. \
Prioritise European/German signals; only include global signals if highly relevant."""


def _signals_to_text(extracted, categories, url_map):
    lines = []
    ref_num = 1
    ref_map = {}
    for cat in categories:
        for item in extracted.get(cat, []):
            for src_idx in item.get("sources", []):
                url = url_map.get((cat, src_idx), "")
                if url and (cat, src_idx) not in ref_map:
                    ref_map[(cat, src_idx)] = ref_num
                    ref_num += 1
    for cat in categories:
        items = extracted.get(cat, [])
        if not items:
            continue
        lines.append(f"### {cat}")
        for item in items:
            refs = [ref_map[(cat, s)] for s in item.get("sources", []) if (cat, s) in ref_map]
            ref_str = " ".join(f"[{r}]" for r in refs)
            lines.append(f"- {item['signal']} {ref_str}")
        lines.append("")
    return "\n".join(lines), ref_map


def _resolve_urls(sources_list, inv_ref_map, url_map):
    return [
        url_map.get(inv_ref_map[s], "")
        for s in sources_list
        if s in inv_ref_map and url_map.get(inv_ref_map[s], "")
    ]


def synthesise_report(extracted, url_map, area, products, competitors, users, categories):
    system = SYNTH_SYSTEM_TEMPLATE.format(
        area=area, products=products, competitors=competitors, users=users
    )
    signals_text, ref_map = _signals_to_text(extracted, categories, url_map)
    user_prompt = f"Market signals from European EV charging industry:\n\n{signals_text}"

    try:
        data = call_groq(system, user_prompt, max_tokens=4000)
    except Exception as e:
        print(f"  synthesis error for {area}: {e}", file=sys.stderr)
        data = {k: [] for k in ["trends", "competitor_moves", "unmet_needs",
                                  "active_regulations", "strategic_implications", "open_questions"]}

    inv_ref_map = {v: k for k, v in ref_map.items()}

    def res(item, *keys):
        return _resolve_urls(item.get("sources", []), inv_ref_map, url_map)

    trends = []
    for t in data.get("trends", [])[:5]:
        trends.append({
            "title": t.get("title", ""),
            "signal": t.get("signal", ""),
            "meaning": t.get("meaning", ""),
            "implication": t.get("implication", ""),
            "urls": res(t),
        })

    competitor_moves = []
    for x in data.get("competitor_moves", []):
        competitor_moves.append({
            "company": x.get("company", ""),
            "action": x.get("action", ""),
            "detail": x.get("detail", x.get("text", "")),
            "urls": res(x),
        })

    unmet_needs = []
    for x in data.get("unmet_needs", []):
        unmet_needs.append({
            "need": x.get("need", x.get("text", "")),
            "evidence": x.get("evidence", ""),
            "urls": res(x),
        })

    regulations = []
    for x in data.get("active_regulations", []):
        regulations.append({
            "regulation": x.get("regulation", x.get("text", "")),
            "status": x.get("status", ""),
            "impact": x.get("impact", ""),
            "urls": res(x),
        })

    implications = []
    for x in data.get("strategic_implications", []):
        implications.append({
            "implication": x.get("implication", x.get("text", "")),
            "rationale": x.get("rationale", ""),
            "urls": [],
        })

    open_questions = [
        {"text": x.get("text", ""), "urls": []}
        for x in data.get("open_questions", [])
    ]

    return {
        "trends": trends,
        "competitor_moves": competitor_moves,
        "unmet_needs": unmet_needs,
        "active_regulations": regulations,
        "strategic_implications": implications,
        "open_questions": open_questions,
    }


# ---------------------------------------------------------------------------
# Phase 4a — Programmatic PPTX generation (all trends on one slide)
# ---------------------------------------------------------------------------

def _tb(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _slide_header(slide, title, period):
    """Clear placeholders then draw title + green accent line + period subtitle."""
    for ph in slide.placeholders:
        try:
            ph.text = ""
        except Exception:
            pass
    _tb(slide, title,
        Inches(0.35), Inches(0.18), Inches(12.6), Inches(0.55),
        size=20, bold=True, color=WHITE)
    line = slide.shapes.add_shape(1, Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = ELLI_GREEN; line.line.width = 0
    _tb(slide, period,
        Inches(0.35), Inches(0.76), Inches(12.6), Inches(0.28),
        size=9, color=ELLI_GREEN)


def draw_trends_slide(slide, trends, area_label, period):
    _slide_header(slide, f"{area_label} — Top 5 Market Trends", period)

    ROW_H   = Inches(1.18)
    START_Y = Inches(1.12)
    NUM_W   = Inches(0.44)
    GAP     = Inches(0.14)
    TEXT_X  = Inches(0.35) + NUM_W + GAP
    TEXT_W  = Inches(12.55) - NUM_W - GAP

    for i, trend in enumerate(trends[:5]):
        y = START_Y + i * ROW_H

        badge = slide.shapes.add_shape(5, Inches(0.35), y + Inches(0.22), NUM_W, NUM_W)
        badge.fill.solid(); badge.fill.fore_color.rgb = ELLI_GREEN; badge.line.width = 0
        btf = badge.text_frame; btf.word_wrap = False
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = str(i + 1)
        br.font.bold = True; br.font.size = Pt(13); br.font.color.rgb = TEXT_ON_GREEN

        _tb(slide, trend.get("title", ""),
            TEXT_X, y + Inches(0.10), TEXT_W, Inches(0.38),
            size=11, bold=True, color=WHITE, wrap=False)

        signal  = (trend.get("signal", "") or "").strip()
        impl    = (trend.get("implication", "") or "").strip()
        desc = signal
        if impl and len(desc) < 220:
            desc = desc.rstrip(".") + ". " + impl if desc else impl
        desc = (desc or (trend.get("meaning", "") or "").strip())[:320]

        _tb(slide, desc,
            TEXT_X, y + Inches(0.48), TEXT_W, Inches(0.62),
            size=8.5, color=LIGHT_GRAY, wrap=True)

        if i < 4:
            sep = slide.shapes.add_shape(1,
                Inches(0.35), y + ROW_H - Inches(0.04),
                Inches(12.6), Inches(0.015))
            sep.fill.solid(); sep.fill.fore_color.rgb = SEPARATOR; sep.line.width = 0


def draw_list_slide(slide, title, items, period, item_fn):
    _slide_header(slide, title, period)
    y = Inches(1.10)
    for item in items:
        used = item_fn(slide, item, y)
        y += used + Inches(0.12)
        if y > Inches(7.0):
            break


def _competitor_item(slide, item, y):
    company = item.get("company", "")
    action  = item.get("action", "")
    detail  = item.get("detail", "")
    label   = f"{company}: {action}" if company else action
    _tb(slide, f"▸ {label}",
        Inches(0.35), y, Inches(12.6), Inches(0.35),
        size=10, bold=True, color=ELLI_GREEN)
    if detail:
        _tb(slide, detail,
            Inches(0.55), y + Inches(0.35), Inches(12.4), Inches(0.55),
            size=9, color=LIGHT_GRAY, wrap=True)
        return Inches(0.9)
    return Inches(0.35)


def _need_item(slide, item, y):
    need     = item.get("need", "")
    evidence = item.get("evidence", "")
    _tb(slide, f"● {need}",
        Inches(0.35), y, Inches(12.6), Inches(0.35),
        size=10, bold=True, color=ELLI_GREEN)
    if evidence:
        _tb(slide, evidence,
            Inches(0.55), y + Inches(0.35), Inches(12.4), Inches(0.55),
            size=9, color=LIGHT_GRAY, wrap=True)
        return Inches(0.9)
    return Inches(0.35)


def _reg_item(slide, item, y):
    reg    = item.get("regulation", "")
    status = item.get("status", "")
    impact = item.get("impact", "")
    header = f"▸ {reg}" + (f"  [{status}]" if status else "")
    _tb(slide, header,
        Inches(0.35), y, Inches(12.6), Inches(0.35),
        size=10, bold=True, color=ELLI_GREEN)
    if impact:
        _tb(slide, impact,
            Inches(0.55), y + Inches(0.35), Inches(12.4), Inches(0.6),
            size=9, color=LIGHT_GRAY, wrap=True)
        return Inches(0.95)
    return Inches(0.35)


def _impl_item(slide, item, y):
    impl      = item.get("implication", "")
    rationale = item.get("rationale", "")
    _tb(slide, f"→ {impl}",
        Inches(0.35), y, Inches(12.6), Inches(0.35),
        size=10, bold=True, color=ELLI_GREEN)
    if rationale:
        _tb(slide, rationale,
            Inches(0.55), y + Inches(0.35), Inches(12.4), Inches(0.6),
            size=9, color=LIGHT_GRAY, wrap=True)
        return Inches(0.95)
    return Inches(0.35)


def _question_item(slide, item, y):
    _tb(slide, f"? {item.get('text', '')}",
        Inches(0.35), y, Inches(12.6), Inches(0.38),
        size=10, color=WHITE, wrap=True)
    return Inches(0.38)


def _clear_and_set_placeholder(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return


def _rich_body(slide, idx, sections):
    """
    Fill body placeholder (idx) with richly formatted sections.
    sections = list of (label, text, label_color) where label may be None.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != idx:
            continue
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True

        first = True
        for label, text, label_color in sections:
            if not text:
                continue
            if label:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                if not first:
                    p.space_before = Pt(8)
                r = p.add_run()
                r.text = label
                r.font.bold = True
                r.font.size = Pt(10)
                r.font.color.rgb = label_color
                first = False

            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            r2 = p2.add_run()
            r2.text = text
            r2.font.size = Pt(9)
            r2.font.color.rgb = WHITE
            first = False
        return


def fill_trend_slide(slide, trend, num, total, area_label):
    _clear_and_set_placeholder(slide, 0, trend["title"])
    _clear_and_set_placeholder(slide, 4, f"Trend {num} of {total}  ·  {area_label}")
    _rich_body(slide, 15, [
        ("SIGNAL", trend.get("signal", ""), ELLI_GREEN),
        ("WHAT IT MEANS", trend.get("meaning", ""), ELLI_GREEN),
        ("ELLI IMPLICATION", trend.get("implication", ""), ELLI_GREEN),
    ])


def fill_competitor_slide(slide, items, subtitle):
    _clear_and_set_placeholder(slide, 4, subtitle)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 15:
            continue
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                p.space_before = Pt(10)
            r = p.add_run()
            company = item.get("company", "")
            action = item.get("action", "")
            r.text = f"{company}: {action}" if company else action
            r.font.bold = True
            r.font.size = Pt(10)
            r.font.color.rgb = ELLI_GREEN
            first = False

            detail = item.get("detail", "")
            if detail:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                r2 = p2.add_run()
                r2.text = detail
                r2.font.size = Pt(9)
                r2.font.color.rgb = WHITE
        return


def fill_needs_slide(slide, items, subtitle):
    _clear_and_set_placeholder(slide, 4, subtitle)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 15:
            continue
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                p.space_before = Pt(10)
            r = p.add_run()
            r.text = f"• {item.get('need', item.get('text', ''))}"
            r.font.bold = True
            r.font.size = Pt(10)
            r.font.color.rgb = ELLI_GREEN
            first = False

            evidence = item.get("evidence", "")
            if evidence:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                r2 = p2.add_run()
                r2.text = evidence
                r2.font.size = Pt(9)
                r2.font.color.rgb = WHITE
        return


def fill_regulations_slide(slide, items, subtitle):
    _clear_and_set_placeholder(slide, 4, subtitle)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 15:
            continue
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                p.space_before = Pt(10)
            r = p.add_run()
            reg = item.get("regulation", item.get("text", ""))
            status = item.get("status", "")
            r.text = f"▸ {reg}" + (f"  [{status}]" if status else "")
            r.font.bold = True
            r.font.size = Pt(10)
            r.font.color.rgb = ELLI_GREEN
            first = False

            impact = item.get("impact", "")
            if impact:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                r2 = p2.add_run()
                r2.text = impact
                r2.font.size = Pt(9)
                r2.font.color.rgb = WHITE
        return


def fill_implications_slide(slide, items, subtitle):
    _clear_and_set_placeholder(slide, 4, subtitle)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 15:
            continue
        tf = ph.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                p.space_before = Pt(10)
            r = p.add_run()
            r.text = f"→ {item.get('implication', item.get('text', ''))}"
            r.font.bold = True
            r.font.size = Pt(10)
            r.font.color.rgb = ELLI_GREEN
            first = False

            rationale = item.get("rationale", "")
            if rationale:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                r2 = p2.add_run()
                r2.text = rationale
                r2.font.size = Pt(9)
                r2.font.color.rgb = WHITE
        return


def fill_pptx(template_path, output_path, fleet_report, site_report, dates):
    """Fill the 14-slide template with programmatically drawn content."""
    prs    = Presentation(template_path)
    slides = list(prs.slides)
    period = f"{dates['period_start']} – {dates['period_end']}"

    # Slide 0: Cover
    for shape in slides[0].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = (run.text
                        .replace("{{GENERATED_AT}}", dates["generated_at"])
                        .replace("{{PERIOD_START}}", dates["period_start"])
                        .replace("{{PERIOD_END}}",   dates["period_end"]))

    # Slide 2: Fleet Trends (all 5 on one slide)
    draw_trends_slide(slides[2], fleet_report["trends"],
                      "Fleet Mobility Management", period)

    # Slide 3: Fleet Competitor Moves
    draw_list_slide(slides[3], "Fleet Mobility — Competitor Moves",
                    fleet_report["competitor_moves"], period, _competitor_item)

    # Slide 4: Fleet Unmet Customer Needs
    draw_list_slide(slides[4], "Fleet Mobility — Unmet Customer Needs",
                    fleet_report["unmet_needs"], period, _need_item)

    # Slide 5: Fleet Active Regulations
    draw_list_slide(slides[5], "Fleet Mobility — Active Regulations",
                    fleet_report["active_regulations"], period, _reg_item)

    # Slide 6: Fleet Strategic Implications
    draw_list_slide(slides[6], "Fleet Mobility — Strategic Implications for Elli",
                    fleet_report["strategic_implications"], period, _impl_item)

    # Slide 8: Site Trends (all 5 on one slide)
    draw_trends_slide(slides[8], site_report["trends"],
                      "Charging Site Management", period)

    # Slide 9: Site Competitor Moves
    draw_list_slide(slides[9], "Charging Site — Competitor Moves",
                    site_report["competitor_moves"], period, _competitor_item)

    # Slide 10: Site Unmet Customer Needs
    draw_list_slide(slides[10], "Charging Site — Unmet Customer Needs",
                    site_report["unmet_needs"], period, _need_item)

    # Slide 11: Site Active Regulations
    draw_list_slide(slides[11], "Charging Site — Active Regulations",
                    site_report["active_regulations"], period, _reg_item)

    # Slide 12: Site Strategic Implications
    draw_list_slide(slides[12], "Charging Site — Strategic Implications for Elli",
                    site_report["strategic_implications"], period, _impl_item)

    # Slide 13: Open Questions
    draw_list_slide(slides[13], "Open Questions for Follow-up",
                    fleet_report["open_questions"] + site_report["open_questions"],
                    period, _question_item)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Phase 4b — Generate Markdown
# ---------------------------------------------------------------------------

def _cite(urls):
    if not urls:
        return ""
    return " " + " ".join(f"[[{i+1}]]({u})" for i, u in enumerate(urls) if u)


def _sources_section(report, area_label):
    seen = set()
    sources = []
    for section in ["trends", "competitor_moves", "unmet_needs",
                     "active_regulations", "strategic_implications"]:
        for item in report.get(section, []):
            for url in item.get("urls", []):
                if url and url not in seen:
                    seen.add(url)
                    sources.append(url)
    if not sources:
        return ""
    lines = [f"### {area_label}"]
    for i, url in enumerate(sources, 1):
        lines.append(f"{i}. {url}")
    return "\n".join(lines) + "\n"


def generate_markdown(fleet_report, site_report, dates):
    label = dates["generated_at"]
    period = f"{dates['period_start']} – {dates['period_end']}"
    today = dates["end"].strftime("%Y-%m-%d")

    lines = [
        f"# Elli Market Signals — {label}",
        f"*Research period: {period} | Generated: {today} | Geographic focus: Germany & Europe*",
        "",
        "---",
        "",
        "## Fleet Mobility Management",
        "",
        "### Top 5 Trends",
        "",
    ]

    for i, t in enumerate(fleet_report["trends"], 1):
        lines.append(f"#### {i}. {t['title']}{_cite(t['urls'])}")
        if t.get("signal"):
            lines.append(f"**Signal:** {t['signal']}")
        if t.get("meaning"):
            lines.append(f"**What it means:** {t['meaning']}")
        if t.get("implication"):
            lines.append(f"**Elli implication:** {t['implication']}")
        lines.append("")

    lines += ["### Competitor Moves", ""]
    for x in fleet_report["competitor_moves"]:
        company = x.get("company", "")
        action = x.get("action", "")
        detail = x.get("detail", "")
        header = f"**{company}:** {action}" if company else f"**{action}**"
        lines.append(f"- {header}{_cite(x['urls'])}")
        if detail:
            lines.append(f"  {detail}")

    lines += ["", "### Unmet Customer Needs", ""]
    for x in fleet_report["unmet_needs"]:
        lines.append(f"- **{x.get('need', '')}**{_cite(x['urls'])}")
        if x.get("evidence"):
            lines.append(f"  {x['evidence']}")

    lines += ["", "### Active Regulations", ""]
    for x in fleet_report["active_regulations"]:
        status = f" [{x['status']}]" if x.get("status") else ""
        lines.append(f"- **{x.get('regulation', '')}**{status}{_cite(x['urls'])}")
        if x.get("impact"):
            lines.append(f"  {x['impact']}")

    lines += ["", "### Strategic Implications for Elli", ""]
    for x in fleet_report["strategic_implications"]:
        lines.append(f"- **{x.get('implication', '')}**")
        if x.get("rationale"):
            lines.append(f"  {x['rationale']}")

    lines += ["", "---", "", "## Charging Site Management", "", "### Top 5 Trends", ""]

    for i, t in enumerate(site_report["trends"], 1):
        lines.append(f"#### {i}. {t['title']}{_cite(t['urls'])}")
        if t.get("signal"):
            lines.append(f"**Signal:** {t['signal']}")
        if t.get("meaning"):
            lines.append(f"**What it means:** {t['meaning']}")
        if t.get("implication"):
            lines.append(f"**Elli implication:** {t['implication']}")
        lines.append("")

    lines += ["### Competitor Moves", ""]
    for x in site_report["competitor_moves"]:
        company = x.get("company", "")
        action = x.get("action", "")
        detail = x.get("detail", "")
        header = f"**{company}:** {action}" if company else f"**{action}**"
        lines.append(f"- {header}{_cite(x['urls'])}")
        if detail:
            lines.append(f"  {detail}")

    lines += ["", "### Unmet Customer Needs", ""]
    for x in site_report["unmet_needs"]:
        lines.append(f"- **{x.get('need', '')}**{_cite(x['urls'])}")
        if x.get("evidence"):
            lines.append(f"  {x['evidence']}")

    lines += ["", "### Active Regulations", ""]
    for x in site_report["active_regulations"]:
        status = f" [{x['status']}]" if x.get("status") else ""
        lines.append(f"- **{x.get('regulation', '')}**{status}{_cite(x['urls'])}")
        if x.get("impact"):
            lines.append(f"  {x['impact']}")

    lines += ["", "### Strategic Implications for Elli", ""]
    for x in site_report["strategic_implications"]:
        lines.append(f"- **{x.get('implication', '')}**")
        if x.get("rationale"):
            lines.append(f"  {x['rationale']}")

    lines += ["", "---", "", "## Open Questions", ""]
    for q in fleet_report["open_questions"] + site_report["open_questions"]:
        lines.append(f"- {q.get('text', '')}")

    lines += ["", "---", "", "## Sources", "",
              _sources_section(fleet_report, "Fleet Mobility Management"),
              "",
              _sources_section(site_report, "Charging Site Management")]

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Phase 5 — Commit MD + post to Slack
# ---------------------------------------------------------------------------

def commit_markdown(md_path, label):
    repo_root = Path(__file__).parent
    token = os.getenv("GITHUB_TOKEN", "")
    try:
        subprocess.run(["git", "config", "user.email", "pipeline@elli-market-signals"], check=True, cwd=repo_root)
        subprocess.run(["git", "config", "user.name", "Market Signals Pipeline"], check=True, cwd=repo_root)
        subprocess.run(["git", "add", str(md_path)], check=True, cwd=repo_root)
        subprocess.run(["git", "commit", "-m", f"report: market signals {label}"], check=True, cwd=repo_root)
        result = subprocess.run(["git", "remote", "get-url", "origin"],
                                 capture_output=True, text=True, cwd=repo_root)
        remote_url = result.stdout.strip()
        if token and "github.com" in remote_url and "@" not in remote_url:
            auth_url = remote_url.replace("https://", f"https://x-access-token:{token}@")
            subprocess.run(["git", "push", auth_url, "HEAD"], check=True, cwd=repo_root)
        else:
            subprocess.run(["git", "push", "-u", "origin", "HEAD"], check=True, cwd=repo_root)
        clean_url = remote_url.replace(".git", "")
        branch = subprocess.run(["git", "branch", "--show-current"],
                                  capture_output=True, text=True, cwd=repo_root).stdout.strip()
        return f"{clean_url}/blob/{branch}/reports/{md_path.name}"
    except subprocess.CalledProcessError as e:
        print(f"  git error: {e}", file=sys.stderr)
        return ""


def post_to_slack(pptx_path, md_path, github_url):
    token = os.environ["SLACK_BOT_TOKEN"]
    channel = os.environ["SLACK_CHANNEL_ID"]
    headers = {"Authorization": f"Bearer {token}"}
    message = ":bar_chart: *Elli Market Signals report is ready!*\n"
    if github_url:
        message += f":link: <{github_url}|View full cited report on GitHub>"
    requests.post("https://slack.com/api/chat.postMessage",
                  headers=headers, json={"channel": channel, "text": message}).raise_for_status()
    for fpath, fname in [(pptx_path, "Market Signals Report.pptx"), (md_path, md_path.name)]:
        if not fpath.exists():
            continue
        with open(fpath, "rb") as f:
            resp = requests.post("https://slack.com/api/files.upload", headers=headers,
                                  data={"channels": channel, "filename": fname},
                                  files={"file": (fname, f)})
        resp.raise_for_status()
        if not resp.json().get("ok"):
            print(f"  Slack upload error for {fname}: {resp.json()}", file=sys.stderr)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    dry_run = os.getenv("DRY_RUN", "false").lower() == "true"
    dates = compute_date_range()
    label = dates["label"]

    print(f"=== Market Signals Pipeline — {dates['generated_at']} ===")
    print(f"Research window: {dates['period_start']} to {dates['period_end']}")
    print(f"DRY_RUN={dry_run}")

    print("\n[Phase 1] Running searches...")
    print("  Fleet categories:")
    fleet_raw = run_searches(FLEET_CATEGORIES, FLEET_COMPETITOR_QUERIES, dates["month_range"])
    print("  Site categories:")
    site_raw = run_searches(SITE_CATEGORIES, SITE_COMPETITOR_QUERIES, dates["month_range"])

    fleet_url_map = build_url_map(fleet_raw)
    site_url_map = build_url_map(site_raw)

    print("\n[Phase 2] Extracting signals...")
    print("  Fleet:")
    fleet_extracted = extract_signals(fleet_raw, batch_size=4)
    print("  Site:")
    site_extracted = extract_signals(site_raw, batch_size=4)

    print("\n[Phase 3] Synthesising report...")
    print("  Fleet report...")
    fleet_report = synthesise_report(
        fleet_extracted, fleet_url_map,
        area="Fleet Mobility Management",
        products=FLEET_PRODUCTS, competitors=FLEET_COMPETITORS, users=FLEET_USERS,
        categories=[c for c, _ in FLEET_CATEGORIES],
    )
    print("  Site report...")
    site_report = synthesise_report(
        site_extracted, site_url_map,
        area="Charging Site Management",
        products=SITE_PRODUCTS, competitors=SITE_COMPETITORS, users=SITE_USERS,
        categories=[c for c, _ in SITE_CATEGORIES],
    )

    print("\n[Phase 4] Generating outputs...")
    pptx_path = Path(f"output_{label}.pptx")
    md_path = Path("reports") / f"output_{label}.md"
    md_path.parent.mkdir(exist_ok=True)

    fill_pptx("template.pptx", str(pptx_path), fleet_report, site_report, dates)
    print(f"  ✓ PPTX: {pptx_path}")

    md_content = generate_markdown(fleet_report, site_report, dates)
    md_path.write_text(md_content, encoding="utf-8")
    print(f"  ✓ MD: {md_path}")

    if dry_run:
        print(f"\n[DRY RUN] Artifacts: {pptx_path}, {md_path}")
        return

    print("\n[Phase 5] Committing MD and posting to Slack...")
    github_url = commit_markdown(md_path, label)
    print(f"  GitHub URL: {github_url or '(push failed)'}")
    post_to_slack(pptx_path, md_path, github_url)
    print("  ✓ Posted to Slack")
    print(f"\n=== Done — {dates['generated_at']} ===")


if __name__ == "__main__":
    main()
