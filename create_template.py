"""
One-time script: generates template.pptx (22 slides) from Elli corporate template.
Run: python create_template.py
Then commit the resulting template.pptx to the repo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SOURCE = "Elli_Mobility_PP_Template_2.0.pptx"
OUTPUT = "template.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xD4, 0x8A)


def remove_all_slides(prs):
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        slide_id_list.remove(sld_id)


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def set_ph(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return


def add_content_slide(prs, layout, title, body_placeholder, subtitle=None):
    slide = prs.slides.add_slide(layout)
    set_ph(slide, 0, title)
    set_ph(slide, 15, body_placeholder)
    if subtitle is not None:
        try:
            set_ph(slide, 4, subtitle)
        except Exception:
            pass
    return slide


def add_section_divider(prs, layout, section_title, section_number):
    """Dark section break slide matching the Elli source template style."""
    slide = prs.slides.add_slide(layout)
    # Large section title (72pt bold white, same position as source slide 3)
    tb = slide.shapes.add_textbox(
        Inches(1.92), Inches(2.29), Inches(9.52), Inches(2.41)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = section_title
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # Section number (large, left side)
    nb = slide.shapes.add_textbox(
        Inches(0.57), Inches(2.06), Inches(1.40), Inches(2.86)
    )
    nf = nb.text_frame
    para2 = nf.paragraphs[0]
    run2 = para2.add_run()
    run2.text = section_number
    run2.font.size = Pt(120)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    return slide


def main():
    prs = Presentation(SOURCE)
    remove_all_slides(prs)

    dark_content = get_layout(prs, "Dark background content")
    dark_empty = get_layout(prs, "Dark background empty")
    chapter = get_layout(prs, "6_Chapter Page / Break")
    light_content = get_layout(prs, "Light background content")

    # ── Slide 0: Cover ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(chapter)
    set_ph(slide, 0, "Market Signals — {{GENERATED_AT}}")
    set_ph(slide, 15, "Research period: {{PERIOD_START}} – {{PERIOD_END}}")

    # ── Slide 1: Fleet section divider ─────────────────────────────────────
    add_section_divider(prs, dark_empty, "Fleet Mobility\nManagement", "1")

    # ── Slides 2–6: Fleet Trends (one per trend) ───────────────────────────
    for n in range(1, 6):
        add_content_slide(
            prs, dark_content,
            title=f"Fleet — Trend {n} of 5",
            subtitle=f"Trend {n} of 5  ·  Fleet Mobility Management",
            body_placeholder="[Signal, meaning and implication will be filled by pipeline]",
        )

    # ── Slide 7: Fleet Competitor Moves ────────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Fleet Mobility — Competitor Moves",
        subtitle="",
        body_placeholder="[Competitor moves will be filled by pipeline]",
    )

    # ── Slide 8: Fleet Unmet Customer Needs ────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Fleet Mobility — Unmet Customer Needs",
        subtitle="",
        body_placeholder="[Unmet needs will be filled by pipeline]",
    )

    # ── Slide 9: Fleet Active Regulations ──────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Fleet Mobility — Active Regulations",
        subtitle="",
        body_placeholder="[Regulations will be filled by pipeline]",
    )

    # ── Slide 10: Fleet Strategic Implications ─────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Fleet Mobility — Strategic Implications for Elli",
        subtitle="",
        body_placeholder="[Strategic implications will be filled by pipeline]",
    )

    # ── Slide 11: Site section divider ─────────────────────────────────────
    add_section_divider(prs, dark_empty, "Charging Site\nManagement", "2")

    # ── Slides 12–16: Site Trends (one per trend) ──────────────────────────
    for n in range(1, 6):
        add_content_slide(
            prs, dark_content,
            title=f"Site — Trend {n} of 5",
            subtitle=f"Trend {n} of 5  ·  Charging Site Management",
            body_placeholder="[Signal, meaning and implication will be filled by pipeline]",
        )

    # ── Slide 17: Site Competitor Moves ────────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Charging Site — Competitor Moves",
        subtitle="",
        body_placeholder="[Competitor moves will be filled by pipeline]",
    )

    # ── Slide 18: Site Unmet Customer Needs ────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Charging Site — Unmet Customer Needs",
        subtitle="",
        body_placeholder="[Unmet needs will be filled by pipeline]",
    )

    # ── Slide 19: Site Active Regulations ──────────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Charging Site — Active Regulations",
        subtitle="",
        body_placeholder="[Regulations will be filled by pipeline]",
    )

    # ── Slide 20: Site Strategic Implications ──────────────────────────────
    add_content_slide(
        prs, dark_content,
        title="Charging Site — Strategic Implications for Elli",
        subtitle="",
        body_placeholder="[Strategic implications will be filled by pipeline]",
    )

    # ── Slide 21: Open Questions ────────────────────────────────────────────
    add_content_slide(
        prs, light_content,
        title="Open Questions for Follow-up",
        subtitle="",
        body_placeholder="[Open questions will be filled by pipeline]",
    )

    prs.save(OUTPUT)
    print(f"✓ Saved {OUTPUT} with {len(prs.slides)} slides")
    print("Slides:")
    for i, slide in enumerate(prs.slides):
        texts = [
            s.text_frame.text[:60].replace('\n', ' ')
            for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        print(f"  {i:2d}: {' | '.join(texts[:2])}")


if __name__ == "__main__":
    main()
